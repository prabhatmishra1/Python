from pptx import Presentation

# Fixing the method for adding bullet points and paragraphs to the slide content
prs = Presentation()

# Slide 1: Title Slide
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
title_1 = slide_1.shapes.title
subtitle_1 = slide_1.placeholders[1]

title_1.text = "Web Programming Seminar"
subtitle_1.text = "Presented by [Your Name]"

# Slide 2: Topics to be Covered
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
title_2 = slide_2.shapes.title
title_2.text = "Topics to be Covered"

content_2 = slide_2.shapes.placeholders[1].text_frame
content_2.text = "We will cover the following topics in this seminar:"
p2 = content_2.add_paragraph()
p2.text = "1. Networking"
p2.level = 0

p2 = content_2.add_paragraph()
p2.text = "2. Frontend"
p2.level = 0

p2 = content_2.add_paragraph()
p2.text = "3. Backend"
p2.level = 0

p2 = content_2.add_paragraph()
p2.text = "4. Database"
p2.level = 0

p2 = content_2.add_paragraph()
p2.text = "5. API"
p2.level = 0

p2 = content_2.add_paragraph()
p2.text = "6. WebSocket"
p2.level = 0

# Slide 3: What Happens When You Enter a URL?
slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
title_3 = slide_3.shapes.title
title_3.text = "What Happens When You Enter a URL?"

content_3 = slide_3.shapes.placeholders[1].text_frame
content_3.text = ("When you enter a URL in the browser, the following steps occur:\n"
                  "1. DNS Lookup: Resolving the domain to an IP address.\n"
                  "2. HTTP/HTTPS Request: The browser sends a request to the server.\n"
                  "3. Backend Processing: The server processes the request and generates a response.\n"
                  "4. Database Query (if necessary): The server retrieves data from the database.\n"
                  "5. Server Response: The server sends back the requested webpage.\n"
                  "6. Rendering: The browser renders the HTML, CSS, and JavaScript.\n"
                  "7. API Calls: Additional data is fetched via APIs if necessary.\n"
                  "8. Caching: Resources are cached for future requests.")

# Slide 4: Networking
slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
title_4 = slide_4.shapes.title
title_4.text = "1. Networking"

content_4 = slide_4.shapes.placeholders[1].text_frame
content_4.text = ("Networking refers to the process of sending data across the internet:\n"
                  "- DNS Lookup translates a domain name into an IP address.\n"
                  "- TCP/IP protocols are used to route data between the browser and server.\n"
                  "- Secure connections use SSL/TLS certificates for HTTPS communication.")

# Slide 5: Frontend
slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
title_5 = slide_5.shapes.title
title_5.text = "2. Frontend"

content_5 = slide_5.shapes.placeholders[1].text_frame
content_5.text = ("The frontend is the visual part of a website that users interact with:\n"
                  "- Technologies like HTML, CSS, and JavaScript are used to build the user interface.\n"
                  "- JavaScript frameworks like React or Vue make the interface dynamic and responsive.\n"
                  "- Browser rendering engines display the website content on the screen.")

# Slide 6: Backend
slide_6 = prs.slides.add_slide(prs.slide_layouts[1])
title_6 = slide_6.shapes.title
title_6.text = "3. Backend"

content_6 = slide_6.shapes.placeholders[1].text_frame
content_6.text = ("The backend manages server-side logic and database interactions:\n"
                  "- It processes HTTP requests and generates responses (HTML, JSON, etc.).\n"
                  "- Backend frameworks like Django, Flask, or Node.js are used to build this logic.\n"
                  "- Ensures security, authentication, and data processing.")

# Slide 7: Database
slide_7 = prs.slides.add_slide(prs.slide_layouts[1])
title_7 = slide_7.shapes.title
title_7.text = "4. Database"

content_7 = slide_7.shapes.placeholders[1].text_frame
content_7.text = ("Databases store and manage data for web applications:\n"
                  "- Common databases include MySQL, PostgreSQL, and MongoDB.\n"
                  "- SQL (Structured Query Language) is used to query relational databases.\n"
                  "- Non-relational databases like MongoDB use document-based storage for scalability.")

# Slide 8: API (Application Programming Interface)
slide_8 = prs.slides.add_slide(prs.slide_layouts[1])
title_8 = slide_8.shapes.title
title_8.text = "5. API"

content_8 = slide_8.shapes.placeholders[1].text_frame
content_8.text = ("APIs allow communication between different parts of an application:\n"
                  "- APIs enable the frontend to interact with the backend without reloading the page.\n"
                  "- REST and GraphQL are common API architectures.\n"
                  "- They facilitate data exchange and allow third-party integrations.")

# Slide 9: WebSocket
slide_9 = prs.slides.add_slide(prs.slide_layouts[1])
title_9 = slide_9.shapes.title
title_9.text = "6. WebSocket"

content_9 = slide_9.shapes.placeholders[1].text_frame
content_9.text = ("WebSockets provide full-duplex communication between the server and client:\n"
                  "- Unlike HTTP, WebSockets maintain an open connection for real-time data transfer.\n"
                  "- Commonly used in chat applications, live updates, and stock trading platforms.\n"
                  "- Allows efficient communication without the overhead of HTTP requests.")

# Slide 10: Conclusion
slide_10 = prs.slides.add_slide(prs.slide_layouts[1])
title_10 = slide_10.shapes.title
title_10.text = "Conclusion"

content_10 = slide_10.shapes.placeholders[1].text_frame
content_10.text = ("Web programming involves multiple technologies working together:\n"
                   "- Networking connects browsers to servers.\n"
                   "- The frontend manages user interaction, while the backend processes logic.\n"
                   "- Databases store data, and APIs/WebSockets enable real-time communication.\n\n"
                   "Thank you for your attention! Any questions?")

# Save the presentation
pptx_file = "Web_Programming_Seminar_Presentation_v2.pptx"
prs.save(pptx_file)

pptx_file
